"""Generate ClaimFlow Presentation v2 - AI Strategist Demo."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
NAVY = RGBColor(30, 58, 95)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
LIGHT_GRAY = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
TEXT_P = RGBColor(15, 23, 42)
TEXT_S = RGBColor(100, 116, 139)
LIGHT_GREEN = RGBColor(236, 253, 245)
LIGHT_TEAL = RGBColor(224, 242, 254)
LIGHT_AMBER = RGBColor(255, 251, 235)
LIGHT_RED = RGBColor(254, 226, 226)
LIGHT_NAVY = RGBColor(226, 232, 240)
BORDER = RGBColor(226, 232, 240)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=WHITE, border=None, bw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=16, color=TEXT_P, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = "Arial"
    p.alignment = align
    return tb


def mtxt(slide, l, t, w, h, lines, default_sz=14, default_color=TEXT_P):
    """Multi-line text with per-line formatting: [(text, size, color, bold), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text = item[0]
        sz = item[1] if len(item) > 1 else default_sz
        color = item[2] if len(item) > 2 else default_color
        bold = item[3] if len(item) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Arial"
        p.space_after = Pt(2)
    return tb


def hline(slide, l, t, w, color=NAVY):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(2))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


# ═══════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

txt(s, Inches(0), Inches(1.2), Inches(10), Inches(1), "ClaimFlow", 72, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(2), Inches(10), Inches(0.6), "AI-Powered Insurance Claims Triage", 28, TEAL, False, PP_ALIGN.CENTER)
hline(s, Inches(2.5), Inches(2.8), Inches(5), WHITE)
txt(s, Inches(0), Inches(3.5), Inches(10), Inches(0.5), "Michael Gardner", 32, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(4), Inches(10), Inches(0.4), "Tenex AI Strategist \u2014 Build First Submission", 20, LIGHT_NAVY, False, PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(4.4), Inches(10), Inches(0.3), "March 2026", 18, LIGHT_NAVY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT_GRAY)

txt(s, Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.45), "The Problem: Manual Triage Slows Everything Down", 34, NAVY, True)

# Volume context
box(s, Inches(0.5), Inches(0.8), Inches(9), Inches(0.45), WHITE, BORDER)
txt(s, Inches(0.7), Inches(0.83), Inches(8.6), Inches(0.35),
    "Typical mid-size insurer: 10,000 claims/month  \u2022  15 adjusters  \u2022  Each handles 650+ claims/year",
    14, TEXT_S, False, PP_ALIGN.CENTER)

# LEFT: Workflow
txt(s, Inches(0.5), Inches(1.4), Inches(5), Inches(0.35), "Current Workflow", 22, NAVY, True)

steps = [
    ("1. Claim Submitted", "15 min \u2014 Policyholder fills 22-field intake form", None, None),
    ("2. Manual Routing & Assignment", "15 min \u2014 Adjuster reads, decides queue; Manager assigns", LIGHT_RED, RED),
    ("3. Manual Workflow Setup", "10 min \u2014 Re-reads form, checks prior claims, fraud patterns", LIGHT_AMBER, AMBER),
    ("4. Investigation", "2 hours \u2014 Actual claim review and decision-making", None, None),
    ("5. Decision Made", "Approved, denied, or escalated", None, None),
]
for i, (step, detail, bg_c, bd_c) in enumerate(steps):
    y = Inches(1.85) + Inches(i * 0.55)
    if bg_c:
        box(s, Inches(0.45), y - Inches(0.03), Inches(5.5), Inches(0.5), bg_c, bd_c, 2)
    txt(s, Inches(0.6), y, Inches(4), Inches(0.22), step, 14, TEXT_P, True)
    txt(s, Inches(0.6), y + Inches(0.2), Inches(5.2), Inches(0.22), detail, 11, TEXT_S)

# Cost labels
txt(s, Inches(4.5), Inches(1.85) + Inches(1 * 0.55), Inches(1.4), Inches(0.2), "$1.8M/yr \u2192", 11, RED, True, PP_ALIGN.RIGHT)
txt(s, Inches(4.5), Inches(1.85) + Inches(2 * 0.55), Inches(1.4), Inches(0.2), "$1.2M/yr \u2192", 11, AMBER, True, PP_ALIGN.RIGHT)

# RIGHT: Cost + Value
txt(s, Inches(6.2), Inches(1.4), Inches(3.4), Inches(0.3), "The Challenge", 18, NAVY, True)
txt(s, Inches(6.2), Inches(1.75), Inches(3.4), Inches(0.6),
    "Before investigation begins, adjusters spend 25 min per claim on triage, assignment, and workflow setup.", 13, TEXT_P)

txt(s, Inches(6.2), Inches(2.4), Inches(3.4), Inches(0.3), "The Cost", 18, NAVY, True)
box(s, Inches(6.2), Inches(2.7), Inches(3.4), Inches(1.1), WHITE, BORDER)
mtxt(s, Inches(6.4), Inches(2.75), Inches(3), Inches(1.0), [
    ("10,000 claims/month", 12, TEXT_P),
    ("\u00d7 25 minutes saved", 12, TEXT_P),
    ("\u00d7 $60/hour adjuster cost", 12, TEXT_P),
    ("\u00d7 12 months", 12, TEXT_P),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", 12, TEXT_S),
    ("= $3,000,000/year", 16, AMBER, True),
])

txt(s, Inches(6.2), Inches(3.9), Inches(3.4), Inches(0.25), "Two Value Propositions", 14, NAVY, True)
box(s, Inches(6.2), Inches(4.15), Inches(3.4), Inches(0.3), LIGHT_GREEN)
txt(s, Inches(6.35), Inches(4.17), Inches(3.1), Inches(0.25), "Efficiency: Same volume, 20% fewer adjusters", 11, TEXT_P)
box(s, Inches(6.2), Inches(4.48), Inches(3.4), Inches(0.3), LIGHT_TEAL)
txt(s, Inches(6.35), Inches(4.5), Inches(3.1), Inches(0.25), "Growth: Scale to 15K claims without hiring", 11, TEXT_P)

# Target market
box(s, Inches(0.5), Inches(4.9), Inches(9), Inches(0.55), WHITE, BORDER)
txt(s, Inches(0.7), Inches(4.92), Inches(1.5), Inches(0.2), "Target Market:", 12, NAVY, True)
txt(s, Inches(0.7), Inches(5.12), Inches(8.6), Inches(0.3),
    "Mid-size auto insurers (10K\u201350K claims/mo)  |  Collision claims <$10K (60% of volume, lowest regulatory risk)  |  $50M\u2013$500M revenue companies",
    11, TEXT_S)

# ═══════════════════════════════════════
# SLIDE 3: WHY I CHOSE THIS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT_GRAY)

txt(s, Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.45), "Why I Chose This Problem", 34, NAVY, True)

# Personal story
box(s, Inches(0.5), Inches(0.85), Inches(9), Inches(0.55), WHITE, BORDER)
txt(s, Inches(0.7), Inches(0.9), Inches(8.6), Inches(0.45),
    "\U0001F697  I got in a fender bender a couple years ago \u2014 the insurance process was a mess. I wanted to fix the low-hanging fruit: the 25 minutes of manual work before investigation even begins.",
    14, TEXT_P)

# Three criteria
criteria = [
    ("\u2713", "Specific", "One workflow step,\nmeasurable impact"),
    ("\u2713", "Solvable", "Deterministic logic,\nno hallucination risk"),
    ("\u2713", "Scalable", "$3M \u2192 $12M\nROI roadmap"),
]
for i, (icon, header, detail) in enumerate(criteria):
    x = Inches(0.5) + Inches(i * 3.1)
    box(s, x, Inches(1.6), Inches(2.8), Inches(0.85), WHITE, BORDER)
    txt(s, x + Inches(0.15), Inches(1.65), Inches(0.4), Inches(0.3), icon, 20, GREEN, True)
    txt(s, x + Inches(0.5), Inches(1.65), Inches(2.1), Inches(0.3), header, 18, NAVY, True)
    txt(s, x + Inches(0.5), Inches(1.95), Inches(2.1), Inches(0.4), detail, 12, TEXT_S)

# How ClaimFlow triages
box(s, Inches(0.5), Inches(2.65), Inches(9), Inches(2.1), WHITE, BORDER)
txt(s, Inches(0.7), Inches(2.7), Inches(5), Inches(0.35), "How ClaimFlow Triages", 22, NAVY, True)

queues = [
    ("Fast-Track Queue:", "<$5K damage, other party fault, police report filed, \u22641 prior claims, no fraud signals", GREEN),
    ("Standard Review:", "$5K\u2013$15K OR policyholder fault OR no police report, \u22643 prior claims, no escalation triggers", AMBER),
    ("Senior Review:", ">$15K OR 2+ fraud signals OR >3 prior claims OR damage exceeds coverage limit", RED),
]
for i, (label, detail, color) in enumerate(queues):
    y = Inches(3.1) + Inches(i * 0.45)
    txt(s, Inches(0.7), y, Inches(2), Inches(0.25), label, 14, color, True)
    txt(s, Inches(2.7), y, Inches(6.5), Inches(0.25), detail, 12, TEXT_P)

# Mission
txt(s, Inches(0.7), Inches(4.45), Inches(8.6), Inches(0.25),
    "ClaimFlow speeds up human decision-making \u2014 it doesn\u2019t replace it. AI triages and recommends. Humans decide.",
    14, TEAL, False, PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════
# SLIDE 4: HOW IT WORKS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT_GRAY)

txt(s, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4), "How ClaimFlow Works", 34, NAVY, True)

workflow = [
    ("1. Claim Submitted", "Policyholder fills 22-field intake form with vehicle, damage, fault, police report, prior claims", "15 min", None, TEXT_S),
    ("2. AI Routing & Assignment", "AI reads all 22 fields, applies triage rubric, routes to correct queue, assigns to least-busy adjuster", "10 sec", TEAL, WHITE),
    ("3. Adjuster Sees Claim", "Claim appears in queue with AI recommendation badge, confidence score, fraud signal count, estimated payout", "Instant", None, TEXT_S),
    ("4. Adjuster Reviews", "Opens modal: AI reasoning, decision factors, all data pre-populated \u2014 no need to re-read intake form", "10 min saved", GREEN, WHITE),
    ("5. Adjuster Decides", "One-click approve, deny (with structured codes), or escalate (with notes). Logged in audit trail.", "5 min saved", None, TEXT_S),
]

for i, (header, detail, badge_text, badge_bg, badge_color) in enumerate(workflow):
    y = Inches(0.7) + Inches(i * 0.85)

    # Left border accent for AI steps
    if badge_bg:
        accent = slide_shapes = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Pt(4), Inches(0.75))
        accent.fill.solid()
        accent.fill.fore_color.rgb = badge_bg
        accent.line.fill.background()

    box(s, Inches(0.55), y, Inches(7.8), Inches(0.75), WHITE, BORDER)
    txt(s, Inches(0.75), y + Inches(0.08), Inches(6), Inches(0.25), header, 16, NAVY, True)
    txt(s, Inches(0.75), y + Inches(0.35), Inches(7.3), Inches(0.35), detail, 12, TEXT_P)

    # Badge
    if badge_bg:
        box(s, Inches(8.5), y + Inches(0.15), Inches(1.1), Inches(0.4), badge_bg)
        txt(s, Inches(8.5), y + Inches(0.18), Inches(1.1), Inches(0.35), badge_text, 12, badge_color, True, PP_ALIGN.CENTER)
    else:
        txt(s, Inches(8.5), y + Inches(0.2), Inches(1.1), Inches(0.3), badge_text, 11, badge_color, False, PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 4:
        txt(s, Inches(4.5), y + Inches(0.72), Inches(1), Inches(0.15), "\u25BC", 10, TEXT_S, False, PP_ALIGN.CENTER)

# Bottom callout
box(s, Inches(0.5), Inches(5), Inches(9), Inches(0.4), TEAL)
txt(s, Inches(0.7), Inches(5.03), Inches(8.6), Inches(0.35),
    "Time Saved: 25 minutes per claim (15 min routing + 10 min workflow enhancement)", 16, WHITE, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 5: THE ROI
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT_GRAY)

txt(s, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4), "The ROI: $3M Today, $12M Tomorrow", 34, NAVY, True)

# LEFT: Current State
txt(s, Inches(0.5), Inches(0.7), Inches(4.4), Inches(0.35), "Current State (What\u2019s Built)", 20, GREEN, True)

txt(s, Inches(0.5), Inches(1.05), Inches(4.4), Inches(0.25), "What ClaimFlow Automates:", 16, NAVY, True)
mtxt(s, Inches(0.5), Inches(1.3), Inches(4.4), Inches(0.7), [
    ("\u2713 AI Routing & Assignment (15 min)", 13, TEXT_P),
    ("\u2713 Workflow Enhancement (10 min)", 13, TEXT_P),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", 10, TEXT_S),
    ("Total: 25 minutes saved per claim", 14, GREEN, True),
])

box(s, Inches(0.5), Inches(2.1), Inches(4.4), Inches(1.15), WHITE, GREEN, 3)
mtxt(s, Inches(0.7), Inches(2.15), Inches(4), Inches(1.05), [
    ("10,000 claims/month", 13, TEXT_P),
    ("\u00d7 25 minutes saved", 13, TEXT_P),
    ("\u00d7 $60/hour adjuster cost", 13, TEXT_P),
    ("\u00d7 12 months", 13, TEXT_P),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", 10, TEXT_S),
    ("= $3,000,000/year", 20, GREEN, True),
])

txt(s, Inches(0.5), Inches(3.35), Inches(4.4), Inches(0.2), "Value Delivery:", 14, NAVY, True)
box(s, Inches(0.5), Inches(3.55), Inches(4.4), Inches(0.25), LIGHT_GREEN)
txt(s, Inches(0.65), Inches(3.57), Inches(4.1), Inches(0.2), "\u2713 Same volume, 20% fewer adjusters", 12, TEXT_P)
box(s, Inches(0.5), Inches(3.83), Inches(4.4), Inches(0.25), LIGHT_GREEN)
txt(s, Inches(0.65), Inches(3.85), Inches(4.1), Inches(0.2), "\u2713 Scale to 15K claims without hiring", 12, TEXT_P)

# RIGHT: Future State
txt(s, Inches(5.1), Inches(0.7), Inches(4.4), Inches(0.35), "Future State (Full Automation)", 20, TEAL, True)

txt(s, Inches(5.1), Inches(1.05), Inches(4.4), Inches(0.25), "What Full Integration Adds:", 16, NAVY, True)
mtxt(s, Inches(5.1), Inches(1.3), Inches(4.4), Inches(1.0), [
    ("\u2713 AI Routing & Assignment (15 min) \u2014 built", 12, GREEN),
    ("\u2713 Workflow Enhancement (10 min) \u2014 built", 12, GREEN),
    ("\U0001F52E Document Review (30 min) \u2014 roadmap", 12, TEXT_S),
    ("\U0001F52E Communication Automation (25 min) \u2014 roadmap", 12, TEXT_S),
    ("\U0001F52E Investigation Assistance (40 min) \u2014 roadmap", 12, TEXT_S),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", 10, TEXT_S),
    ("Total: 2 hours saved per claim", 14, TEAL, True),
])

box(s, Inches(5.1), Inches(2.55), Inches(4.4), Inches(1.15), WHITE, TEAL, 3)
mtxt(s, Inches(5.3), Inches(2.6), Inches(4), Inches(1.05), [
    ("10,000 claims/month", 13, TEXT_P),
    ("\u00d7 1.67 hours saved (conservative)", 13, TEXT_P),
    ("\u00d7 $60/hour adjuster cost", 13, TEXT_P),
    ("\u00d7 12 months", 13, TEXT_P),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", 10, TEXT_S),
    ("= $12,000,000/year", 20, TEAL, True),
])

txt(s, Inches(5.1), Inches(3.8), Inches(4.4), Inches(0.2), "Integration Path:", 14, NAVY, True)
path = ["30-day pilot \u2192 Prove accuracy \u2192 Expand scope"]
box(s, Inches(5.1), Inches(4.0), Inches(4.4), Inches(0.25), LIGHT_TEAL)
txt(s, Inches(5.25), Inches(4.02), Inches(4.1), Inches(0.2), path[0], 12, TEAL, True)

# Assumptions
box(s, Inches(0.5), Inches(4.5), Inches(9), Inches(0.95), WHITE, BORDER)
txt(s, Inches(0.7), Inches(4.55), Inches(2), Inches(0.2), "Key Assumptions:", 13, NAVY, True)
txt(s, Inches(0.7), Inches(4.75), Inches(8.6), Inches(0.65),
    "Claim volume: 10,000/month  |  Adjuster cost: $60/hr fully loaded  |  Time saved: 15 min (routing) + 10 min (workflow) = 25 min\n"
    "Conservative future estimate: 1.67 hrs accounts for complexity variance and edge cases requiring manual review\n"
    "ROI does not include fraud prevention savings, reduced reassignment costs, or improved customer satisfaction",
    11, TEXT_S)

# ═══════════════════════════════════════
# SLIDE 6: STRATEGIC TRADE-OFFS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT_GRAY)

txt(s, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4), "Strategic Trade-Offs", 34, NAVY, True)

# LEFT: Tech Stack
txt(s, Inches(0.5), Inches(0.7), Inches(4), Inches(0.3), "Tech Stack (What I Built)", 20, NAVY, True)

stack = [
    ("\u2713 Deterministic AI Logic", "No LLMs, no hallucinations \u2014 every decision\ncites exact field values for regulatory compliance"),
    ("\u2713 SQLite Database", "Fast prototyping, swappable to Postgres\nin 2 hours for production"),
    ("\u2713 React + Tailwind CSS", "Role-based UX, 7 user profiles\n(Admin + 6 adjusters) with optimized workflows"),
    ("\u2713 Fraud Detection", "4 automated signals with visual prominence\n(red borders, fraud-first sorting)"),
    ("\u2713 Dual ROI Model", "Current ($3M) + future ($12M) with\ntransparent assumptions"),
]
for i, (header, detail) in enumerate(stack):
    y = Inches(1.05) + Inches(i * 0.75)
    txt(s, Inches(0.5), y, Inches(4), Inches(0.22), header, 14, TEXT_P, True)
    txt(s, Inches(0.5), y + Inches(0.22), Inches(4), Inches(0.45), detail, 11, TEXT_S)

# RIGHT: Go-Forward
txt(s, Inches(4.8), Inches(0.7), Inches(4.8), Inches(0.3), "Go-Forward Features", 20, NAVY, True)

features = [
    ("\U0001F4E4", "Multi-Claim Batch Processing", "Upload 100+ claims via CSV, bulk triage"),
    ("\U0001F517", "Claims Management Integration", "API connectors to Guidewire, Duck Creek"),
    ("\U0001F9E0", "Advanced Fraud ML Models", "Learn from adjuster decisions over time"),
    ("\U0001F4F1", "Mobile App for Field Adjusters", "On-site processing with photo capture"),
    ("\U0001F4C4", "Document OCR & Parsing", "Extract data from photos, PDFs, forms"),
]
for i, (icon, header, detail) in enumerate(features):
    y = Inches(1.05) + Inches(i * 0.72)
    box(s, Inches(4.8), y, Inches(4.8), Inches(0.62), WHITE, BORDER)
    txt(s, Inches(4.95), y + Inches(0.05), Inches(0.4), Inches(0.3), icon, 16, TEXT_P)
    txt(s, Inches(5.35), y + Inches(0.05), Inches(4), Inches(0.25), header, 14, NAVY, True)
    txt(s, Inches(5.35), y + Inches(0.3), Inches(4), Inches(0.25), detail, 11, TEXT_S)

# Bottom
txt(s, Inches(0.5), Inches(4.9), Inches(9), Inches(0.4),
    "This demo proves the core value: AI can triage accurately and speed up workflows. Production infrastructure is table stakes.",
    14, TEXT_P, False, PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════
# SLIDE 7: CLOSING / DEMO
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, Inches(0), Inches(1.2), Inches(10), Inches(0.7), "ClaimFlow", 48, NAVY, True, PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(1.9), Inches(10), Inches(0.4), "Strategic scoping. Rapid validation. Business outcomes.", 22, TEAL, False, PP_ALIGN.CENTER)

hline(s, Inches(3), Inches(2.5), Inches(4), NAVY)

txt(s, Inches(0), Inches(2.8), Inches(10), Inches(0.3), "Demo: [LIVE LINK]", 18, TEXT_S, False, PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(3.1), Inches(10), Inches(0.3), "GitHub: github.com/mjgard47/tenex-insurance-claims-triage", 16, TEXT_S, False, PP_ALIGN.CENTER)

hline(s, Inches(3), Inches(3.6), Inches(4), NAVY)

txt(s, Inches(0), Inches(4.1), Inches(10), Inches(0.3),
    "Thanks for watching. I\u2019d love to build the future of work with you.", 16, NAVY, False, PP_ALIGN.CENTER, italic=True)

# Save
out = r"C:\Users\Michael Gardner\Desktop\Tenex - Insurance Claims Project\presentation\ClaimFlow_Presentation_v2.pptx"
prs.save(out)
print("ClaimFlow_Presentation_v2.pptx created successfully!")
print(f"Location: {out}")
print(f"Slides: 7")
print(f"Theme: Navy (#1E3A5F) + Teal (#0D9488)")
print(f"ROI: $3M (current) + $12M (future) with transparent math")
print(f"Ready for Tenex AI Strategist demo submission.")
