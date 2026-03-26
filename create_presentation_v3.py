"""ClaimFlow Presentation v3 - Executive-ready, minimal text, speaker notes for detail."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(30, 58, 95)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
LGRAY = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
TP = RGBColor(15, 23, 42)
TS = RGBColor(100, 116, 139)
LGREEN = RGBColor(236, 253, 245)
LTEAL = RGBColor(224, 242, 254)
LAMBER = RGBColor(255, 251, 235)
LRED = RGBColor(254, 226, 226)
LNAVY = RGBColor(226, 232, 240)
BDR = RGBColor(226, 232, 240)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

def bg(sl, c):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = c

def bx(sl, l, t, w, h, f=WHITE, b=None, bw=1):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = f
    if b: s.line.color.rgb = b; s.line.width = Pt(bw)
    else: s.line.fill.background()
    return s

def tx(sl, l, t, w, h, text, sz=16, c=TP, bold=False, al=PP_ALIGN.LEFT, it=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = bold
    p.font.italic = it; p.font.name = "Arial"; p.alignment = al
    return tb

def hl(sl, l, t, w, c=NAVY):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(2))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def notes(sl, text):
    sl.notes_slide.notes_text_frame.text = text

# ═══════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, NAVY)
tx(sl, Inches(0), Inches(1.3), Inches(10), Inches(1), "ClaimFlow", 72, WHITE, True, PP_ALIGN.CENTER)
tx(sl, Inches(0), Inches(2.1), Inches(10), Inches(0.5), "AI-Powered Insurance Claims Triage", 28, TEAL, False, PP_ALIGN.CENTER)
hl(sl, Inches(3), Inches(2.8), Inches(4), WHITE)
tx(sl, Inches(0), Inches(3.4), Inches(10), Inches(0.5), "Michael Gardner", 32, WHITE, True, PP_ALIGN.CENTER)
tx(sl, Inches(0), Inches(3.9), Inches(10), Inches(0.4), "Tenex AI Strategist \u2014 Build First Submission", 20, LNAVY, False, PP_ALIGN.CENTER)
tx(sl, Inches(0), Inches(4.3), Inches(10), Inches(0.3), "March 2026", 18, LNAVY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LGRAY)
tx(sl, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5), "The Problem", 36, NAVY, True)
tx(sl, Inches(0.5), Inches(0.8), Inches(9), Inches(0.4), "Before investigation begins, adjusters waste 25 minutes per claim on manual triage and assignment.", 18, TS)

# Three stat cards
stats = [
    ("15 min", "wasted per claim\non manual triage", AMBER),
    ("$3M", "annual cost at\n10K claims/month", RED),
    ("10 sec", "ClaimFlow triage\nvs 15 min manual", TEAL),
]
for i, (num, label, color) in enumerate(stats):
    x = Inches(0.5) + Inches(i * 3.1)
    bx(sl, x, Inches(1.5), Inches(2.8), Inches(1.8), WHITE, BDR)
    tx(sl, x + Inches(0.3), Inches(1.7), Inches(2.2), Inches(0.8), num, 48, color, True, PP_ALIGN.CENTER)
    tx(sl, x + Inches(0.3), Inches(2.5), Inches(2.2), Inches(0.6), label, 15, TS, False, PP_ALIGN.CENTER)

# Workflow - simple
tx(sl, Inches(0.5), Inches(3.6), Inches(9), Inches(0.3), "Where the waste happens:", 18, NAVY, True)

bx(sl, Inches(0.5), Inches(4.0), Inches(2.6), Inches(0.45), WHITE, BDR)
tx(sl, Inches(0.6), Inches(4.05), Inches(2.4), Inches(0.35), "1. Claim Submitted", 14, TP)

tx(sl, Inches(3.2), Inches(4.1), Inches(0.3), Inches(0.3), "\u2192", 16, TS, False, PP_ALIGN.CENTER)

bx(sl, Inches(3.5), Inches(3.95), Inches(2.9), Inches(0.55), LRED, RED, 2)
tx(sl, Inches(3.6), Inches(3.97), Inches(2.7), Inches(0.2), "2. Manual Triage & Assignment", 13, RED, True)
tx(sl, Inches(3.6), Inches(4.2), Inches(2.7), Inches(0.2), "25 min of waste", 12, RED)

tx(sl, Inches(6.5), Inches(4.1), Inches(0.3), Inches(0.3), "\u2192", 16, TS, False, PP_ALIGN.CENTER)

bx(sl, Inches(6.8), Inches(4.0), Inches(2.7), Inches(0.45), WHITE, BDR)
tx(sl, Inches(6.9), Inches(4.05), Inches(2.5), Inches(0.35), "3. Investigation & Decision", 14, TP)

# Bottom
bx(sl, Inches(0.5), Inches(4.8), Inches(9), Inches(0.5), LAMBER)
tx(sl, Inches(0.7), Inches(4.85), Inches(8.6), Inches(0.4),
   "This is the step ClaimFlow automates \u2014 routing and assignment, not the investigation itself.", 15, TP, False, PP_ALIGN.CENTER, it=True)

notes(sl, """SPEAKER NOTES - The Problem:
- Typical mid-size insurer: 10,000 claims/month, 15 adjusters
- Each adjuster handles 650+ claims/year
- Manual triage: adjuster reads 22-field form, decides complexity, picks queue
- Manual assignment: manager checks workload, assigns to available adjuster
- Combined waste: 25 min/claim x 10K claims x $60/hr x 12 months = $3M/year
- Collision claims under $10K are high-volume and well-suited for automation
- Two value props: efficiency (same volume, fewer adjusters) or growth (more volume, same team)""")

# ═══════════════════════════════════════
# SLIDE 3: WHY THIS PROBLEM
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LGRAY)
tx(sl, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5), "Why I Chose This Problem", 36, NAVY, True)

bx(sl, Inches(0.5), Inches(1.0), Inches(9), Inches(0.6), WHITE, BDR)
tx(sl, Inches(0.7), Inches(1.05), Inches(8.6), Inches(0.5),
   "I got in a fender bender a couple years ago \u2014 the insurance process was a mess. I wanted to fix the low-hanging fruit that AI can solve today.",
   16, TP)

# Three criteria
for i, (icon, title, sub) in enumerate([
    ("\u2713", "Specific", "One workflow step, measurable impact"),
    ("\u2713", "Solvable", "Deterministic logic, no hallucination risk"),
    ("\u2713", "Scalable", "$3M \u2192 $12M ROI roadmap"),
]):
    x = Inches(0.5) + Inches(i * 3.1)
    bx(sl, x, Inches(1.9), Inches(2.8), Inches(0.9), WHITE, BDR)
    tx(sl, x + Inches(0.2), Inches(1.95), Inches(0.4), Inches(0.3), icon, 20, GREEN, True)
    tx(sl, x + Inches(0.6), Inches(1.95), Inches(2), Inches(0.3), title, 20, NAVY, True)
    tx(sl, x + Inches(0.6), Inches(2.3), Inches(2), Inches(0.4), sub, 13, TS)

# Triage rubric
tx(sl, Inches(0.5), Inches(3.1), Inches(9), Inches(0.35), "How ClaimFlow Routes Claims", 22, NAVY, True)

for i, (queue, rule, color) in enumerate([
    ("Fast-Track", "<$5K, other party fault, police report, \u22641 prior claims", GREEN),
    ("Standard Review", "$5K\u2013$15K, or policyholder fault, or no police report", AMBER),
    ("Senior Review", ">$15K, or fraud signals, or >3 prior claims", RED),
]):
    y = Inches(3.55) + Inches(i * 0.5)
    bx(sl, Inches(0.5), y, Inches(1.8), Inches(0.35), color)
    tx(sl, Inches(0.55), y + Inches(0.03), Inches(1.7), Inches(0.3), queue, 13, WHITE, True, PP_ALIGN.CENTER)
    tx(sl, Inches(2.5), y + Inches(0.03), Inches(7), Inches(0.3), rule, 13, TP)

# Bottom
bx(sl, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4), TEAL)
tx(sl, Inches(0.7), Inches(5.03), Inches(8.6), Inches(0.35),
   "AI triages and recommends. Humans decide.", 16, WHITE, True, PP_ALIGN.CENTER)

notes(sl, """SPEAKER NOTES - Why This Problem:
- Personal experience with insurance claims frustration
- Scoped to collision claims under $10K with no injuries
- High-volume claim type, lowest regulatory risk, fastest ROI
- Deterministic AI = no LLMs, no hallucination risk, every decision cites exact fields
- Triage criteria are rule-based: damage amount, fault, police report, prior claims, fraud signals
- 4 fraud signals: airbags deployed but drivable, high mileage + high damage, no police report on high value, excessive prior claims
- ClaimFlow is decision SUPPORT - adjusters retain full authority""")

# ═══════════════════════════════════════
# SLIDE 4: HOW IT WORKS
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LGRAY)
tx(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.4), "How ClaimFlow Works", 36, NAVY, True)

steps = [
    ("1", "Claim Submitted", "22-field intake form", "15 min", None, TS),
    ("2", "AI Triage & Assignment", "Routes to queue, assigns to least-busy adjuster", "10 sec", TEAL, WHITE),
    ("3", "Adjuster Reviews", "AI reasoning pre-populated \u2014 no re-reading the form", "10 min saved", GREEN, WHITE),
    ("4", "Adjuster Decides", "One-click approve, deny, or escalate with audit trail", "5 min saved", None, TS),
]
for i, (num, title, detail, badge, badge_bg, badge_c) in enumerate(steps):
    y = Inches(0.75) + Inches(i * 1.05)
    if badge_bg:
        accent = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Pt(5), Inches(0.85))
        accent.fill.solid(); accent.fill.fore_color.rgb = badge_bg; accent.line.fill.background()
    bx(sl, Inches(0.55), y, Inches(7.6), Inches(0.85), WHITE, BDR)
    tx(sl, Inches(0.75), y + Inches(0.1), Inches(6.5), Inches(0.35), title, 20, NAVY, True)
    tx(sl, Inches(0.75), y + Inches(0.45), Inches(6.5), Inches(0.3), detail, 14, TS)
    if badge_bg:
        bx(sl, Inches(8.3), y + Inches(0.2), Inches(1.2), Inches(0.4), badge_bg)
        tx(sl, Inches(8.3), y + Inches(0.22), Inches(1.2), Inches(0.35), badge, 13, badge_c, True, PP_ALIGN.CENTER)
    else:
        tx(sl, Inches(8.3), y + Inches(0.25), Inches(1.2), Inches(0.3), badge, 12, badge_c, False, PP_ALIGN.CENTER)
    if i < 3:
        tx(sl, Inches(4.3), y + Inches(0.85), Inches(1), Inches(0.2), "\u25BC", 12, TS, False, PP_ALIGN.CENTER)

bx(sl, Inches(0.5), Inches(4.95), Inches(9), Inches(0.45), TEAL)
tx(sl, Inches(0.7), Inches(4.98), Inches(8.6), Inches(0.4),
   "25 minutes saved per claim \u2014 15 min routing + 10 min workflow", 17, WHITE, True, PP_ALIGN.CENTER)

notes(sl, """SPEAKER NOTES - How ClaimFlow Works:
Step 1: Policyholder fills form with vehicle details, damage, fault, police report, prior claims
Step 2: AI reads all 22 fields, applies deterministic triage rubric, routes to Fast-Track/Standard/Senior queue, auto-assigns to least-busy adjuster via workload balancing
Step 3: Adjuster opens claim modal - sees AI reasoning with exact criteria checks (green checkmarks/red X), fraud signal badges, estimated payout calculation, all decision factors pre-populated
Step 4: Approve (one click), Deny (structured reason codes + notes), or Escalate (with handoff notes to next tier). Every action logged in audit trail with timestamp.
The AI processes in ~1.2 seconds. Manual triage takes 15 minutes. That's a 99% reduction in triage time.""")

# ═══════════════════════════════════════
# SLIDE 5: TARGET MARKET
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LGRAY)
tx(sl, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5), "Target Market & User Roles", 36, NAVY, True)

# Market
for i, (title, sub) in enumerate([
    ("Mid-Size Insurers", "10K\u201350K claims/month"),
    ("Collision Claims <$10K", "High volume, lowest risk"),
    ("$50M\u2013$500M Revenue", "Tenex sweet spot"),
]):
    x = Inches(0.5) + Inches(i * 3.1)
    bx(sl, x, Inches(1.0), Inches(2.8), Inches(0.9), WHITE, BDR)
    tx(sl, x + Inches(0.3), Inches(1.1), Inches(2.2), Inches(0.35), title, 17, NAVY, True)
    tx(sl, x + Inches(0.3), Inches(1.5), Inches(2.2), Inches(0.3), sub, 14, TS)

tx(sl, Inches(0.5), Inches(2.2), Inches(9), Inches(0.35), "Four User Roles", 22, NAVY, True)

roles = [
    ("Junior Adjuster", "Fast-Track Queue", "Simple claims, speed-focused", RGBColor(209, 250, 229), RGBColor(6, 95, 70)),
    ("Standard Adjuster", "Standard Review", "Moderate complexity", RGBColor(254, 243, 199), RGBColor(146, 64, 14)),
    ("Senior Adjuster", "Senior Review / SIU", "Fraud investigation", RGBColor(224, 231, 255), RGBColor(49, 46, 129)),
    ("Admin", "Operations Center", "Platform oversight", RGBColor(219, 234, 254), NAVY),
]
for i, (title, queue, desc, bg_c, text_c) in enumerate(roles):
    x = Inches(0.5) + Inches(i * 2.35)
    bx(sl, x, Inches(2.7), Inches(2.1), Inches(1.3), bg_c)
    tx(sl, x + Inches(0.2), Inches(2.8), Inches(1.7), Inches(0.3), title, 16, text_c, True)
    tx(sl, x + Inches(0.2), Inches(3.15), Inches(1.7), Inches(0.25), queue, 13, TP)
    tx(sl, x + Inches(0.2), Inches(3.45), Inches(1.7), Inches(0.25), desc, 12, TS)

bx(sl, Inches(0.5), Inches(4.3), Inches(9), Inches(0.5), WHITE, BDR)
tx(sl, Inches(0.7), Inches(4.35), Inches(8.6), Inches(0.4),
   "Each role sees only what they need. Adjusters land on their queue. Admin sees the full platform.", 15, TP, False, PP_ALIGN.CENTER, it=True)

notes(sl, """SPEAKER NOTES - Target Market:
- Mid-size auto insurers processing 10K-50K claims/month
- Focus on collision claims under $10K with no injuries - 60% of total claim volume
- Lowest regulatory risk, fastest path to proving AI value
- Tenex sweet spot: $50M-$500M revenue companies needing operational transformation
- Junior Adjusters handle Fast-Track queue (simple, <$5K, clear fault)
- Standard Adjusters handle moderate complexity ($5K-$15K, mixed fault)
- Senior Adjusters/SIU handle high-value claims and fraud investigations
- Admin sees Operations Center with ROI metrics, triage criteria reference, all queues""")

# ═══════════════════════════════════════
# SLIDE 6: ROI - CURRENT STATE
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LGRAY)
tx(sl, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5), "ROI: What\u2019s Built Today", 36, NAVY, True)
tx(sl, Inches(0.5), Inches(0.8), Inches(9), Inches(0.35), "Triage-Only Automation \u2014 25 minutes saved per claim", 18, TS)

# Big number
bx(sl, Inches(0.5), Inches(1.4), Inches(4.2), Inches(2.0), WHITE, GREEN, 3)
tx(sl, Inches(0.7), Inches(1.5), Inches(3.8), Inches(0.3), "Projected Annual Savings", 16, TS)
tx(sl, Inches(0.7), Inches(1.9), Inches(3.8), Inches(0.8), "$3,000,000", 48, GREEN, True)
tx(sl, Inches(0.7), Inches(2.7), Inches(3.8), Inches(0.3), "per year at 10K claims/month", 14, TS)

# Math
bx(sl, Inches(5.2), Inches(1.4), Inches(4.3), Inches(2.0), WHITE, BDR)
tx(sl, Inches(5.4), Inches(1.5), Inches(3.9), Inches(0.3), "The Math", 18, NAVY, True)
tx(sl, Inches(5.4), Inches(1.85), Inches(3.9), Inches(0.25), "10,000 claims/month", 15, TP)
tx(sl, Inches(5.4), Inches(2.1), Inches(3.9), Inches(0.25), "\u00d7 25 minutes saved", 15, TP)
tx(sl, Inches(5.4), Inches(2.35), Inches(3.9), Inches(0.25), "\u00d7 $60/hour adjuster cost", 15, TP)
tx(sl, Inches(5.4), Inches(2.6), Inches(3.9), Inches(0.25), "\u00d7 12 months", 15, TP)
hl(sl, Inches(5.4), Inches(2.85), Inches(3.5), BDR)
tx(sl, Inches(5.4), Inches(2.9), Inches(3.9), Inches(0.3), "= $3,000,000/year", 18, GREEN, True)

# What it automates
tx(sl, Inches(0.5), Inches(3.7), Inches(9), Inches(0.3), "What ClaimFlow Automates Today", 18, NAVY, True)
for i, (item, time) in enumerate([
    ("AI reads 22 fields, applies triage rubric, routes to correct queue", "15 min saved"),
    ("Auto-assigns to least-busy adjuster, pre-populates decision factors", "10 min saved"),
]):
    y = Inches(4.1) + Inches(i * 0.45)
    bx(sl, Inches(0.5), y, Inches(7.5), Inches(0.35), LGREEN)
    tx(sl, Inches(0.7), y + Inches(0.03), Inches(6), Inches(0.3), "\u2713 " + item, 13, TP)
    tx(sl, Inches(7.2), y + Inches(0.03), Inches(1.5), Inches(0.3), time, 13, GREEN, True, PP_ALIGN.RIGHT)

notes(sl, """SPEAKER NOTES - ROI Current State:
- 25 minutes saved = 15 min (routing & queue assignment) + 10 min (workflow enhancement, pre-populated data)
- $60/hour is fully loaded adjuster cost including benefits
- 10,000 claims/month is a mid-size insurer baseline
- Two ways to capture value:
  1. Efficiency: Handle same claim volume with 20% fewer adjusters
  2. Growth: Scale from 10K to 15K claims/month without hiring
- ROI does NOT include: fraud prevention savings, reduced reassignment costs, faster customer payouts, improved satisfaction
- 30-day pilot approach: run ClaimFlow in parallel with existing process, compare accuracy""")

# ═══════════════════════════════════════
# SLIDE 7: ROI - FUTURE STATE
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LGRAY)
tx(sl, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5), "ROI: The Full Roadmap", 36, NAVY, True)
tx(sl, Inches(0.5), Inches(0.8), Inches(9), Inches(0.35), "Triage + Workflow Automation \u2014 2 hours saved per claim", 18, TS)

# Big number
bx(sl, Inches(0.5), Inches(1.4), Inches(4.2), Inches(2.0), WHITE, TEAL, 3)
tx(sl, Inches(0.7), Inches(1.5), Inches(3.8), Inches(0.3), "Projected Annual Savings", 16, TS)
tx(sl, Inches(0.7), Inches(1.9), Inches(3.8), Inches(0.8), "$12,000,000", 48, TEAL, True)
tx(sl, Inches(0.7), Inches(2.7), Inches(3.8), Inches(0.3), "per year at 10K claims/month", 14, TS)

# Roadmap items
bx(sl, Inches(5.2), Inches(1.4), Inches(4.3), Inches(2.0), WHITE, BDR)
tx(sl, Inches(5.4), Inches(1.5), Inches(3.9), Inches(0.3), "What Gets Added", 18, NAVY, True)
items = [
    ("\u2713 AI Routing & Assignment", "15 min \u2014 built", GREEN),
    ("\u2713 Workflow Enhancement", "10 min \u2014 built", GREEN),
    ("Document Review Automation", "30 min \u2014 roadmap", TS),
    ("Communication Automation", "25 min \u2014 roadmap", TS),
    ("Investigation Assistance", "40 min \u2014 roadmap", TS),
]
for i, (item, time, color) in enumerate(items):
    y = Inches(1.85) + Inches(i * 0.28)
    tx(sl, Inches(5.4), y, Inches(2.5), Inches(0.25), item, 12, color, color == GREEN)
    tx(sl, Inches(7.8), y, Inches(1.5), Inches(0.25), time, 11, color, False, PP_ALIGN.RIGHT)

# Integration path
tx(sl, Inches(0.5), Inches(3.7), Inches(9), Inches(0.3), "Integration Path", 18, NAVY, True)
path_steps = ["30-day pilot", "Prove accuracy", "Expand scope", "Full integration"]
for i, step in enumerate(path_steps):
    x = Inches(0.5) + Inches(i * 2.35)
    bx(sl, x, Inches(4.1), Inches(2.1), Inches(0.45), LTEAL if i < 2 else WHITE, TEAL if i < 2 else BDR)
    tx(sl, x + Inches(0.15), Inches(4.13), Inches(1.8), Inches(0.35), step, 14, TEAL if i < 2 else TS, i < 2, PP_ALIGN.CENTER)
    if i < 3:
        tx(sl, x + Inches(2.1), Inches(4.15), Inches(0.25), Inches(0.3), "\u2192", 14, TS, False, PP_ALIGN.CENTER)

notes(sl, """SPEAKER NOTES - ROI Future State:
- Conservative estimate: 1.67 hours (100 min) saved, not the full 2 hours
- Accounts for complexity variance and edge cases requiring manual review
- Document Review: OCR for police reports, repair estimates, photos
- Communication Automation: auto-generated status updates, adjuster handoff notes
- Investigation Assistance: AI surfaces relevant prior claims, similar cases, pattern detection
- Integration path: start with 30-day parallel pilot, measure AI accuracy vs human decisions
- If AI matches human accuracy at 95%+, expand scope incrementally
- $12M does NOT include reduced fraud losses or improved customer retention""")

# ═══════════════════════════════════════
# SLIDE 8: TRADE-OFFS
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LGRAY)
tx(sl, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5), "Strategic Trade-Offs", 36, NAVY, True)

# Left: What I built
tx(sl, Inches(0.5), Inches(1.0), Inches(4.2), Inches(0.35), "What I Built", 22, NAVY, True)
built = [
    "Deterministic AI logic (no hallucinations)",
    "Role-based UX (7 profiles)",
    "Fraud detection (4 signals)",
    "Status tracking & audit trail",
    "Batch CSV processing",
    "Dual ROI model",
]
for i, item in enumerate(built):
    y = Inches(1.45) + Inches(i * 0.38)
    tx(sl, Inches(0.5), y, Inches(4.2), Inches(0.3), "\u2713 " + item, 15, TP)

# Right: Go-forward
tx(sl, Inches(5.3), Inches(1.0), Inches(4.2), Inches(0.35), "Go-Forward Features", 22, NAVY, True)
forward = [
    ("Claims System Integration", "Guidewire, Duck Creek API connectors"),
    ("Advanced Fraud ML", "Learn from adjuster decisions over time"),
    ("Mobile App", "On-site processing with photo capture"),
    ("Document OCR", "Extract data from photos, PDFs, forms"),
]
for i, (title, sub) in enumerate(forward):
    y = Inches(1.45) + Inches(i * 0.65)
    bx(sl, Inches(5.3), y, Inches(4.2), Inches(0.55), WHITE, BDR)
    tx(sl, Inches(5.5), y + Inches(0.05), Inches(3.8), Inches(0.25), title, 15, NAVY, True)
    tx(sl, Inches(5.5), y + Inches(0.28), Inches(3.8), Inches(0.22), sub, 12, TS)

# Bottom
tx(sl, Inches(0.5), Inches(4.6), Inches(9), Inches(0.4),
   "Every trade-off prioritized proving the AI works over production infrastructure.", 15, TP, False, PP_ALIGN.CENTER, it=True)

notes(sl, """SPEAKER NOTES - Strategic Trade-Offs:
What I built in < 24 hours:
- Deterministic AI: rule-based, cites exact fields, zero hallucination risk, regulatory compliant
- 7 user profiles: Admin, Junior 1&2, Standard 1&2, Senior 1&2 - each with tailored UX
- Fraud detection: airbags+drivable inconsistency, high mileage+high damage, no police report on high value, excessive claims history
- Full audit trail: every status change logged with who/when/what
- Batch processing: upload 100+ claims via CSV with sample template
- SQLite for speed - swappable to Postgres in 2 hours

What I intentionally didn't build:
- Authentication (Auth0 = 1 day, spent it on fraud detection instead)
- Production database (SQLite proves it works, Postgres doesn't change the value prop)
- Live deployment (localhost demo proves functionality, cloud doesn't change ROI)""")

# ═══════════════════════════════════════
# SLIDE 9: CLOSING
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
tx(sl, Inches(0), Inches(1.2), Inches(10), Inches(0.7), "ClaimFlow", 48, NAVY, True, PP_ALIGN.CENTER)
tx(sl, Inches(0), Inches(1.9), Inches(10), Inches(0.4), "Strategic scoping. Rapid validation. Business outcomes.", 22, TEAL, False, PP_ALIGN.CENTER)
hl(sl, Inches(3), Inches(2.5), Inches(4), NAVY)
tx(sl, Inches(0), Inches(2.8), Inches(10), Inches(0.3), "Demo: [LIVE LINK]", 18, TS, False, PP_ALIGN.CENTER)
tx(sl, Inches(0), Inches(3.1), Inches(10), Inches(0.3), "GitHub: github.com/mjgard47/tenex-insurance-claims-triage", 16, TS, False, PP_ALIGN.CENTER)
hl(sl, Inches(3), Inches(3.6), Inches(4), NAVY)
tx(sl, Inches(0), Inches(4.1), Inches(10), Inches(0.3),
   "Thanks for watching. I'd love to build the future of work with you.", 16, NAVY, False, PP_ALIGN.CENTER, it=True)

# Save
out = r"C:\Users\Michael Gardner\Desktop\Tenex - Insurance Claims Project\presentation\ClaimFlow_Presentation_v3.pptx"
prs.save(out)
print(f"ClaimFlow_Presentation_v3.pptx created successfully!")
print(f"Location: {out}")
print(f"Slides: 9")
print(f"Theme: Navy + Teal, minimal text, speaker notes for detail")
