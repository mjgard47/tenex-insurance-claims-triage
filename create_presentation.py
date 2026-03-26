"""Generate ClaimFlow presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
NAVY = RGBColor(30, 58, 95)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
LIGHT_GRAY = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
TEXT_PRIMARY = RGBColor(15, 23, 42)
TEXT_SECONDARY = RGBColor(100, 116, 139)
LIGHT_GREEN_BG = RGBColor(236, 253, 245)
LIGHT_TEAL_BG = RGBColor(224, 242, 254)
LIGHT_AMBER_BG = RGBColor(255, 251, 235)
LIGHT_NAVY = RGBColor(226, 232, 240)
RED = RGBColor(220, 38, 38)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()
    return shape

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

add_text(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(1.2),
         "ClaimFlow", 72, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.7), Inches(11.333), Inches(0.7),
         "AI-Powered Insurance Claims Triage", 28, TEAL, False, PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.5), Inches(4.333), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = WHITE
line.line.fill.background()

add_text(slide, Inches(1), Inches(4.2), Inches(11.333), Inches(0.7),
         "Michael Gardner", 32, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
         "Tenex AI Strategist — Build First Submission", 20, LIGHT_NAVY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
         "March 2026", 18, LIGHT_NAVY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
         "The Problem: Manual Triage Wastes Time & Money", 36, NAVY, True)

# Left side - 3 stat boxes
stats = [
    ("⏱️", "15 minutes", "wasted per claim on manual triage"),
    ("💰", "$1.48M", "annual waste at 10K claims/month"),
    ("😤", "60%", "of claims are simple collisions"),
]
for i, (icon, number, label) in enumerate(stats):
    y = Inches(1.4) + Inches(i * 1.8)
    add_box(slide, Inches(0.8), y, Inches(4.5), Inches(1.5), WHITE)
    add_text(slide, Inches(1.1), y + Inches(0.15), Inches(0.6), Inches(0.5), icon, 28, TEXT_PRIMARY)
    add_text(slide, Inches(1.7), y + Inches(0.15), Inches(3.3), Inches(0.7), number, 44, AMBER, True)
    add_text(slide, Inches(1.7), y + Inches(0.9), Inches(3.3), Inches(0.4), label, 15, TEXT_SECONDARY)

# Right side - workflow
add_text(slide, Inches(5.8), Inches(1.3), Inches(6.5), Inches(0.5),
         "Current Workflow", 24, NAVY, True)

steps = [
    ("1. Claim Submitted", "15 min — Policyholder fills form", False),
    ("2. Manual Triage", "15 min — Adjuster reads, categorizes  ← THIS IS THE WASTE", True),
    ("3. Manual Assignment", "10 min — Manager assigns to adjuster", False),
    ("4. Investigation", "2 hours — Adjuster reviews evidence", False),
    ("5. Decision Made", "Approved, denied, or escalated", False),
]
for i, (step, detail, highlight) in enumerate(steps):
    y = Inches(1.9) + Inches(i * 0.85)
    if highlight:
        add_box(slide, Inches(5.7), y - Inches(0.05), Inches(6.8), Inches(0.8), RGBColor(254, 226, 226), RED)
    add_text(slide, Inches(5.9), y, Inches(6.3), Inches(0.35), step, 17, TEXT_PRIMARY, True)
    c = AMBER if highlight else TEXT_SECONDARY
    add_text(slide, Inches(5.9), y + Inches(0.32), Inches(6.3), Inches(0.35), detail, 13, c, highlight)

# Bottom callout
add_box(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8), LIGHT_AMBER_BG)
add_text(slide, Inches(1.1), Inches(6.5), Inches(11.1), Inches(0.6),
         "Before any investigation begins, adjusters waste 15 minutes deciding if a claim is simple or complex.", 17, TEXT_PRIMARY)

# ============================================================
# SLIDE 3: MY STORY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
         "Why I Chose This Problem", 36, NAVY, True)

# Personal story card
add_box(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(2.2), WHITE)
add_text(slide, Inches(1.1), Inches(1.4), Inches(0.6), Inches(0.5), "🚗", 28, TEXT_PRIMARY)
add_text(slide, Inches(1.7), Inches(1.4), Inches(4), Inches(0.4), "Personal Experience", 22, NAVY, True)
story = ("A couple years ago, I got in a fender bender. The insurance experience was a mess:\n"
         "• Lost my claim twice\n"
         "• Re-submitted documentation 3 times\n"
         "• Took 6 weeks to get paid\n\n"
         "When I started thinking about problems to solve, I wanted low-hanging fruit that AI could solve today.")
add_text(slide, Inches(1.7), Inches(1.9), Inches(10.5), Inches(1.5), story, 15, TEXT_PRIMARY)

# Three criteria boxes
criteria = [
    ("✓", "Specific", "One workflow step, measurable impact"),
    ("✓", "Solvable", "Deterministic logic, no hallucination risk"),
    ("✓", "Scalable", "$1.48M → $12M ROI roadmap"),
]
for i, (icon, header, text) in enumerate(criteria):
    x = Inches(0.8) + Inches(i * 4.0)
    add_box(slide, x, Inches(3.8), Inches(3.7), Inches(1.3), WHITE)
    add_text(slide, x + Inches(0.2), Inches(3.9), Inches(0.5), Inches(0.4), icon, 22, GREEN, True)
    add_text(slide, x + Inches(0.7), Inches(3.9), Inches(2.8), Inches(0.4), header, 20, NAVY, True)
    add_text(slide, x + Inches(0.7), Inches(4.35), Inches(2.8), Inches(0.5), text, 14, TEXT_SECONDARY)

# Bottom callout
add_box(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.0), TEAL)
add_text(slide, Inches(1.1), Inches(5.6), Inches(11.1), Inches(0.8),
         "I scoped to collision claims under $10K with no injuries — that's 60% of claim volume, lowest regulatory risk, fastest ROI.",
         19, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: BEFORE/AFTER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
         "ClaimFlow: Automated Triage in 10 Seconds", 36, NAVY, True)

# LEFT - Current State
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "Current State", 24, AMBER, True)
current_steps = [
    ("Claim Submitted", "⏱️ 15 min", False),
    ("Manual Triage", "⏱️ 15 min", True),
    ("Manual Assignment", "⏱️ 10 min", False),
    ("Investigation", "⏱️ 2 hours", False),
    ("Decision Made", "", False),
]
for i, (step, time, highlight) in enumerate(current_steps):
    y = Inches(1.9) + Inches(i * 0.7)
    if highlight:
        add_box(slide, Inches(0.7), y - Inches(0.05), Inches(5.5), Inches(0.65), RGBColor(254, 226, 226), RED)
    add_text(slide, Inches(1.0), y, Inches(3.5), Inches(0.35), step, 16, TEXT_PRIMARY, True)
    if time:
        add_text(slide, Inches(4.5), y, Inches(1.5), Inches(0.35), time, 14, TEXT_SECONDARY)

add_box(slide, Inches(0.8), Inches(5.5), Inches(5.3), Inches(0.6), LIGHT_AMBER_BG)
add_text(slide, Inches(1.0), Inches(5.55), Inches(5.0), Inches(0.45), "Total: 2 hours 40 min", 18, AMBER, True, PP_ALIGN.CENTER)

# RIGHT - ClaimFlow State
add_text(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.5), "ClaimFlow State", 24, GREEN, True)
cf_steps = [
    ("Claim Submitted", "⏱️ 15 min", False),
    ("AI Triage", "⏱️ 10 seconds", True),
    ("Auto-Assignment", "⏱️ 5 seconds", True),
    ("Investigation", "⏱️ 2 hours", False),
    ("Decision Made", "", False),
]
for i, (step, time, highlight) in enumerate(cf_steps):
    y = Inches(1.9) + Inches(i * 0.7)
    if highlight:
        add_box(slide, Inches(6.7), y - Inches(0.05), Inches(5.8), Inches(0.65), LIGHT_GREEN_BG, GREEN)
    add_text(slide, Inches(7.0), y, Inches(3.5), Inches(0.35), step, 16, TEXT_PRIMARY, True)
    if time:
        add_text(slide, Inches(10.5), y, Inches(1.8), Inches(0.35), time, 14, TEXT_SECONDARY)

add_box(slide, Inches(6.8), Inches(5.5), Inches(5.5), Inches(0.6), LIGHT_GREEN_BG)
add_text(slide, Inches(7.0), Inches(5.55), Inches(5.2), Inches(0.45), "Total: 2 hours 15 min  ·  15 min saved/claim", 18, GREEN, True, PP_ALIGN.CENTER)

# Bottom ROI boxes
add_box(slide, Inches(0.8), Inches(6.3), Inches(5.8), Inches(0.9), LIGHT_GREEN_BG)
add_text(slide, Inches(1.0), Inches(6.35), Inches(3), Inches(0.35), "Triage-Only ROI", 16, GREEN, True)
add_text(slide, Inches(4.0), Inches(6.35), Inches(2.3), Inches(0.5), "$1.48M/year", 32, GREEN, True)
add_text(slide, Inches(1.0), Inches(6.8), Inches(5.3), Inches(0.3), "30-day pilot, parallel running", 13, TEXT_SECONDARY)

add_box(slide, Inches(7.0), Inches(6.3), Inches(5.5), Inches(0.9), LIGHT_TEAL_BG)
add_text(slide, Inches(7.2), Inches(6.35), Inches(3), Inches(0.35), "Full Workflow ROI", 16, TEAL, True)
add_text(slide, Inches(10.2), Inches(6.35), Inches(2.1), Inches(0.5), "$12M/year", 32, TEAL, True)
add_text(slide, Inches(7.2), Inches(6.8), Inches(5.0), Inches(0.3), "Full integration, 2hr saved/claim", 13, TEXT_SECONDARY)

# ============================================================
# SLIDE 5: TARGET MARKET
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
         "Strategic Scoping: 60% of Volume, Lowest Risk", 36, NAVY, True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5), "Target Market", 24, NAVY, True)

market = [
    ("🏢", "Mid-Size Insurers", "10K–50K claims/month", "Regional & national carriers"),
    ("🚗", "Collision Claims", "Under $10K damage", "No injuries (lower risk)"),
    ("📊", "60% of Volume", "High-frequency work", "Fast ROI validation"),
]
for i, (icon, header, line1, line2) in enumerate(market):
    x = Inches(0.8) + Inches(i * 4.0)
    add_box(slide, x, Inches(1.7), Inches(3.7), Inches(1.5), WHITE)
    add_text(slide, x + Inches(0.2), Inches(1.8), Inches(0.5), Inches(0.4), icon, 22, TEXT_PRIMARY)
    add_text(slide, x + Inches(0.7), Inches(1.8), Inches(2.8), Inches(0.4), header, 19, NAVY, True)
    add_text(slide, x + Inches(0.7), Inches(2.25), Inches(2.8), Inches(0.35), line1, 15, TEXT_PRIMARY)
    add_text(slide, x + Inches(0.7), Inches(2.6), Inches(2.8), Inches(0.35), line2, 13, TEXT_SECONDARY)

add_text(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.5), "User Roles", 24, NAVY, True)

roles = [
    (RGBColor(209, 250, 229), "Junior Adjuster", "Fast-Track", "Speed-focused", RGBColor(6, 95, 70)),
    (RGBColor(254, 243, 199), "Standard Adjuster", "Standard Review", "Context + escalation", RGBColor(146, 64, 14)),
    (RGBColor(224, 231, 255), "Senior Adjuster", "Fraud Investigation", "Red flags surfaced", RGBColor(49, 46, 129)),
    (RGBColor(219, 234, 254), "Admin", "Operations Center", "Workload monitoring", NAVY),
]
for i, (bg, title, subtitle, desc, text_color) in enumerate(roles):
    x = Inches(0.8) + Inches(i * 3.1)
    add_box(slide, x, Inches(4.1), Inches(2.8), Inches(1.6), bg)
    add_text(slide, x + Inches(0.2), Inches(4.2), Inches(2.4), Inches(0.4), title, 16, text_color, True)
    add_text(slide, x + Inches(0.2), Inches(4.6), Inches(2.4), Inches(0.35), subtitle, 14, TEXT_PRIMARY)
    add_text(slide, x + Inches(0.2), Inches(4.95), Inches(2.4), Inches(0.3), desc, 12, TEXT_SECONDARY)

# Bottom callout
add_box(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.9), TEAL)
add_text(slide, Inches(1.1), Inches(6.2), Inches(11.1), Inches(0.7),
         "Each role sees only what they need. Admin sees the full platform. Adjusters see their queue and personal stats.",
         18, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: TECH CHOICES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
         "Built in a Day: Strategic Trade-Offs", 36, NAVY, True)

# Left - Tech Stack
add_text(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5), "Tech Stack", 24, NAVY, True)

stack = [
    ("✓ Deterministic AI Logic", "No LLMs, no hallucinations. Every decision cites exact fields."),
    ("✓ SQLite Database", "Fast prototyping. Swappable to Postgres in 2 hours."),
    ("✓ React + Tailwind", "Role-based UX built in hours. 7 user profiles."),
    ("✓ Fraud Detection", "4 automated signals with visual prominence."),
    ("✓ Dual ROI Model", "Phase 1 vs Phase 3 with clear assumptions."),
]
for i, (header, detail) in enumerate(stack):
    y = Inches(1.9) + Inches(i * 0.85)
    add_text(slide, Inches(0.8), y, Inches(5), Inches(0.35), header, 16, TEXT_PRIMARY, True)
    add_text(slide, Inches(0.8), y + Inches(0.3), Inches(5), Inches(0.4), detail, 13, TEXT_SECONDARY)

# Right - Trade-offs
add_text(slide, Inches(6.5), Inches(1.3), Inches(6), Inches(0.5), "Strategic Trade-Offs", 24, NAVY, True)

tradeoffs = [
    ("❌ No Authentication", "Auth0 = 1 day. Spent it on fraud detection instead.", "Add in Phase 2"),
    ("❌ No Production Database", "SQLite proves it works. Postgres doesn't change the value prop.", "Swap to Postgres in 2 hours"),
    ("❌ No Live Deployment", "Localhost demo proves functionality. AWS doesn't change ROI.", "Deploy to AWS in Phase 2"),
]
for i, (header, detail, future) in enumerate(tradeoffs):
    y = Inches(1.9) + Inches(i * 1.4)
    add_box(slide, Inches(6.5), y, Inches(6), Inches(1.2), WHITE)
    add_text(slide, Inches(6.7), y + Inches(0.1), Inches(5.6), Inches(0.35), header, 17, TEXT_PRIMARY, True)
    add_text(slide, Inches(6.7), y + Inches(0.45), Inches(5.6), Inches(0.35), detail, 13, TEXT_SECONDARY)
    add_text(slide, Inches(6.7), y + Inches(0.8), Inches(5.6), Inches(0.3), future, 12, TEAL)

# Bottom callout
add_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9), TEAL)
add_text(slide, Inches(1.1), Inches(6.4), Inches(11.1), Inches(0.7),
         "Every trade-off prioritized proving the AI works over production infrastructure. That's what pilots are for.",
         19, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(1),
         "ClaimFlow", 48, NAVY, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(0.6),
         "Built in an afternoon and evening", 24, TEAL, False, PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.3), Inches(4.333), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = NAVY
line.line.fill.background()

add_text(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(0.6),
         "Strategic scoping. Rapid validation. Business outcomes.", 20, TEXT_PRIMARY, False, PP_ALIGN.CENTER)

line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.4), Inches(4.333), Pt(2))
line2.fill.solid()
line2.fill.fore_color.rgb = NAVY
line2.line.fill.background()

add_text(slide, Inches(1), Inches(4.8), Inches(11.333), Inches(0.4),
         "GitHub: github.com/mjgard47/tenex-insurance-claims-triage", 16, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.4),
         "Demo: localhost:5173 (setup in README)", 16, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
         "Thanks for watching. I'd love to build the future of work with you.", 16, NAVY, False, PP_ALIGN.CENTER)

# Save
output_path = r"C:\Users\Michael Gardner\Desktop\Tenex - Insurance Claims Project\presentation\ClaimFlow_Presentation.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
